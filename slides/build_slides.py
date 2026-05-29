"""Build the three course decks (proposal, interim, final) as BOTH .pptx and .pdf.

Why this approach: there is no LibreOffice on the build machine, so we render the
SAME structured slide spec twice — once via python-pptx (editable .pptx) and once
via matplotlib's PdfPages (.pdf). Decks are defined as plain data in
slides/content/{proposal,interim,final}.py and may pull live numbers from
results/eda_stats.json and ml/results/eval_report.json when present.

Run:  python slides/build_slides.py            # builds all three
      python slides/build_slides.py proposal   # build one
"""
from __future__ import annotations

import importlib
import sys
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

NAVY = (16, 36, 62)
BLUE = (43, 108, 176)
GREY = (74, 85, 104)
ACCENT = (221, 107, 32)

# 16:9
PPT_W, PPT_H = Inches(13.333), Inches(7.5)
PDF_W, PDF_H = 13.333, 7.5


# --------------------------------------------------------------------------- #
# PPTX renderer
# --------------------------------------------------------------------------- #
def _rgb(t):
    return RGBColor(*t)


def _ppt_title_slide(prs, slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _rgb(NAVY)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.4)).text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = slide["title"]
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = _rgb((255, 255, 255))
    if slide.get("subtitle"):
        sp = tb.add_paragraph(); sp.text = slide["subtitle"]
        sp.font.size = Pt(20); sp.font.color.rgb = _rgb((200, 210, 225))
    return s


def _ppt_content_slide(prs, slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # title bar
    bar = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9)).text_frame
    bar.word_wrap = True
    p = bar.paragraphs[0]; p.text = slide["title"]
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = _rgb(NAVY)

    has_img = bool(slide.get("image"))
    body_w = Inches(7.1) if has_img else Inches(12.3)

    if slide.get("bullets"):
        tf = s.shapes.add_textbox(Inches(0.6), Inches(1.4), body_w, Inches(5.4)).text_frame
        tf.word_wrap = True
        first = True
        for b in slide["bullets"]:
            lvl, txt = b if isinstance(b, tuple) else (0, b)
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = ("• " if lvl == 0 else "– ") + txt
            p.level = lvl
            p.font.size = Pt(16 if lvl == 0 else 13)
            p.font.color.rgb = _rgb(NAVY if lvl == 0 else GREY)
            p.space_after = Pt(6)

    if slide.get("table"):
        _ppt_table(s, slide["table"], top=Inches(4.4) if slide.get("bullets") else Inches(1.6))

    if has_img:
        img = ROOT / slide["image"]
        if img.exists():
            s.shapes.add_picture(str(img), Inches(7.5), Inches(1.4), width=Inches(5.4))

    if slide.get("note"):
        nt = s.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.2), Inches(0.5)).text_frame
        nt.word_wrap = True
        p = nt.paragraphs[0]; p.text = slide["note"]
        p.font.size = Pt(11); p.font.italic = True; p.font.color.rgb = _rgb(GREY)
    return s


def _ppt_table(slide, table, top):
    headers, rows = table["headers"], table["rows"]
    nr, nc = len(rows) + 1, len(headers)
    gt = slide.shapes.add_table(nr, nc, Inches(0.6), top, Inches(12.1), Inches(0.4 * nr)).table
    for j, h in enumerate(headers):
        c = gt.cell(0, j); c.text = str(h)
        c.text_frame.paragraphs[0].font.size = Pt(12)
        c.text_frame.paragraphs[0].font.bold = True
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = str(val)
            c.text_frame.paragraphs[0].font.size = Pt(11)


def build_pptx(deck, out: Path):
    prs = Presentation()
    prs.slide_width = PPT_W; prs.slide_height = PPT_H
    for slide in deck["slides"]:
        if slide.get("kind") == "title":
            _ppt_title_slide(prs, slide)
        else:
            _ppt_content_slide(prs, slide)
    prs.save(str(out))


# --------------------------------------------------------------------------- #
# PDF renderer (matplotlib)
# --------------------------------------------------------------------------- #
def _c(t):
    return tuple(v / 255 for v in t)


def _pdf_slide(pdf, slide):
    fig = plt.figure(figsize=(PDF_W, PDF_H))
    ax = fig.add_axes([0, 0, 1, 1]); ax.axis("off")
    ax.set_xlim(0, 100); ax.set_ylim(0, 100)

    if slide.get("kind") == "title":
        ax.add_patch(plt.Rectangle((0, 0), 100, 100, color=_c(NAVY)))
        ax.text(8, 52, slide["title"], fontsize=30, fontweight="bold", color="white", va="center", wrap=True)
        if slide.get("subtitle"):
            ax.text(8, 42, slide["subtitle"], fontsize=15, color=_c((200, 210, 225)), va="center")
        pdf.savefig(fig); plt.close(fig); return

    ax.text(4, 92, slide["title"], fontsize=22, fontweight="bold", color=_c(NAVY), va="center")
    ax.add_patch(plt.Rectangle((4, 88), 92, 0.6, color=_c(BLUE)))

    has_img = bool(slide.get("image")) and (ROOT / slide["image"]).exists()
    # matplotlib's wrap=True wraps to the figure, not a column, so wrap manually
    import textwrap
    base_chars = 58 if has_img else 104

    y = 82
    for b in slide.get("bullets", []):
        lvl, txt = b if isinstance(b, tuple) else (0, b)
        prefix = "•  " if lvl == 0 else "–  "
        indent = 5 + lvl * 4
        lines = textwrap.wrap(txt, width=max(20, base_chars - lvl * 8)) or [""]
        for k, line in enumerate(lines):
            ax.text(indent, y, (prefix if k == 0 else "   ") + line,
                    fontsize=13 if lvl == 0 else 11,
                    color=_c(NAVY if lvl == 0 else GREY), va="top")
            y -= 4.4 if lvl == 0 else 3.8
        y -= 1.6  # gap between bullets

    if slide.get("table"):
        _pdf_table(ax, slide["table"], y0=(y - 2) if slide.get("bullets") else 80)

    if has_img:
        import matplotlib.image as mpimg
        img = mpimg.imread(str(ROOT / slide["image"]))
        ax_img = fig.add_axes([0.55, 0.12, 0.42, 0.66]); ax_img.axis("off")
        ax_img.imshow(img)

    if slide.get("note"):
        ax.text(5, 5, slide["note"], fontsize=10, style="italic", color=_c(GREY), va="center")
    pdf.savefig(fig); plt.close(fig)


def _pdf_table(ax, table, y0):
    headers, rows = table["headers"], table["rows"]
    full = [headers] + [[str(c) for c in r] for r in rows]
    nc = len(headers)
    col_w = 92 / nc
    row_h = min(5.0, (y0 - 8) / max(len(full), 1))
    for i, row in enumerate(full):
        yy = y0 - i * row_h
        for j, cell in enumerate(row):
            ax.text(5 + j * col_w, yy, cell, fontsize=9.5 if i else 10,
                    fontweight="bold" if i == 0 else "normal",
                    color=_c(NAVY), va="top", wrap=True)
        ax.add_patch(plt.Rectangle((4.5, yy - row_h + 1.2), 92, 0.15, color=_c((200, 200, 200))))


def build_pdf(deck, out: Path):
    with PdfPages(str(out)) as pdf:
        for slide in deck["slides"]:
            _pdf_slide(pdf, slide)


# --------------------------------------------------------------------------- #
def build(name: str):
    mod = importlib.import_module(f"content.{name}")
    deck = mod.build()
    build_pptx(deck, HERE / f"{name}.pptx")
    build_pdf(deck, HERE / f"{name}.pdf")
    print(f"  {name}: {len(deck['slides'])} slides → {name}.pptx + {name}.pdf")


if __name__ == "__main__":
    names = sys.argv[1:] or ["proposal", "interim", "final"]
    print("Building decks…")
    for n in names:
        build(n)
